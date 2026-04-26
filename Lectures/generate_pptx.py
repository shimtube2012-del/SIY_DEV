#!/usr/bin/env python3
"""
tech-evolution-slides.html → PowerPoint (.pptx) 변환 스크립트
디자인(어두운 배경, 강조색), 형식(테이블, 카드, 피라미드, 플로우), 내용을 최대한 유지
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── 색상 정의 ──
BG_DARK   = RGBColor(0x0a, 0x0a, 0x1a)
BG_MID    = RGBColor(0x1a, 0x1a, 0x2e)
TEXT_MAIN  = RGBColor(0xe8, 0xe8, 0xf0)
TEXT_SUB   = RGBColor(0xcc, 0xd6, 0xf6)
TEXT_DIM   = RGBColor(0x88, 0x92, 0xb0)
TEXT_FAINT = RGBColor(0x60, 0x6a, 0x86)
ACCENT_RED    = RGBColor(0xe9, 0x45, 0x60)
ACCENT_CYAN   = RGBColor(0x00, 0xd2, 0xff)
ACCENT_GREEN  = RGBColor(0x00, 0xe6, 0x76)
ACCENT_GOLD   = RGBColor(0xff, 0xd7, 0x00)
ACCENT_PURPLE = RGBColor(0x7c, 0x4d, 0xff)
WHITE = RGBColor(0xff, 0xff, 0xff)
CARD_BG = RGBColor(0x14, 0x14, 0x28)
CARD_BORDER = RGBColor(0x30, 0x30, 0x50)
TABLE_HEADER_BG = RGBColor(0x2a, 0x15, 0x20)
TABLE_ROW_BG = RGBColor(0x10, 0x10, 0x22)
TABLE_ROW_ALT = RGBColor(0x14, 0x14, 0x2a)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

FONT_TITLE = 'Malgun Gothic'
FONT_BODY  = 'Malgun Gothic'

# ── 유틸리티 ──

def add_slide():
    """빈 슬라이드 추가 + 어두운 배경 설정"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = BG_DARK
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = BG_MID
    fill.gradient_stops[1].position = 1.0
    return slide


def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox.text_frame


def set_run(run, text, size=18, bold=False, italic=False, color=TEXT_SUB, font_name=FONT_BODY):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def add_text(tf, text, size=18, bold=False, italic=False, color=TEXT_SUB, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.add_paragraph() if len(tf.paragraphs) > 0 and tf.paragraphs[0].text != '' else tf.paragraphs[0]
    if tf.paragraphs[0].text != '' or len(tf.paragraphs) > 1:
        p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    set_run(run, text, size, bold, italic, color)
    return p


def add_para(tf, size=18, bold=False, color=TEXT_SUB, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_run_to_para(p, text, size=18, bold=False, italic=False, color=TEXT_SUB):
    run = p.add_run()
    set_run(run, text, size, bold, italic, color)
    return run


def title_text(slide, text, top=Inches(0.4), size=36, color=TEXT_MAIN, alignment=PP_ALIGN.CENTER):
    tf = add_textbox(slide, Inches(0.5), top, SW - Inches(1), Inches(1))
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    set_run(run, text, size, True, color=color)
    return tf


def subtitle_text(slide, text, top=Inches(1.3), size=20, color=TEXT_DIM, alignment=PP_ALIGN.CENTER):
    tf = add_textbox(slide, Inches(0.5), top, SW - Inches(1), Inches(0.6))
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    set_run(run, text, size, False, color=color)
    return tf


def part_marker(slide, text, top=Inches(0.2)):
    tf = add_textbox(slide, Inches(0.5), top, SW - Inches(1), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, 14, False, color=TEXT_FAINT)


def bridge_text(slide, text, top=Inches(0.15)):
    tf = add_textbox(slide, Inches(1), top, SW - Inches(2), Inches(0.6))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run(run, text, 13, False, True, TEXT_DIM)
    # left border simulation
    line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), top, Inches(0.06), Inches(0.5))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = ACCENT_RED
    line_shape.line.fill.background()


def bridge_end(slide, text, top=Inches(6.8)):
    tf = add_textbox(slide, Inches(1), top, SW - Inches(2), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    set_run(run, text, 12, False, True, TEXT_FAINT)


def source_footer(slide, text, top=Inches(7.0)):
    tf = add_textbox(slide, Inches(0.8), top, SW - Inches(1.6), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run(run, text, 9, False, False, TEXT_FAINT)


def icon_box(slide, emoji, left, top, size=Inches(1.2), color=ACCENT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x10, 0x18)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = emoji
    run.font.size = Pt(48)
    tf.paragraphs[0].space_before = Pt(6)
    return shape


def add_card(slide, left, top, width, height, icon, title, desc, border_color=ACCENT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.size = Pt(28)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    set_run(run2, title, 13, True, color=TEXT_MAIN)
    if desc:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(3)
        run3 = p3.add_run()
        set_run(run3, desc, 10, False, color=TEXT_DIM)
    return shape


def add_card_left_align(slide, left, top, width, height, icon_title, desc, border_color=ACCENT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run(run, icon_title, 14, True, color=TEXT_MAIN)
    if desc:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        set_run(run2, desc, 11, False, color=TEXT_DIM)
    return shape


def add_tag(slide, text, left, top, color=ACCENT_RED, width=None):
    w = width or Inches(1.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(
        min(255, color[0] // 5 + 0x0a),
        min(255, color[1] // 5 + 0x0a),
        min(255, color[2] // 5 + 0x0a)
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, 11, False, color=color)
    return shape


def add_quote_box(slide, text, left, top, width, height, border_color=ACCENT_RED, text_size=20):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x12, 0x0a, 0x14)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(20)
    tf.margin_top = Pt(14)
    tf.margin_bottom = Pt(14)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(text_size * 1.6)
    # split by lines
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            para = p
        else:
            para = tf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            para.line_spacing = Pt(text_size * 1.6)
        # handle **bold** markers
        parts = line.split('**')
        for j, part in enumerate(parts):
            if part == '':
                continue
            run = para.add_run()
            if j % 2 == 1:  # bold
                set_run(run, part, text_size, True, color=ACCENT_RED)
            else:
                set_run(run, part, text_size, False, color=TEXT_MAIN)
    # left border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    border.fill.solid()
    border.fill.fore_color.rgb = border_color
    border.line.fill.background()
    return shape


def add_flow_arrow(slide, left, top):
    tf = add_textbox(slide, left, top, Inches(0.4), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, '→', 18, False, color=TEXT_FAINT)


def add_flow_step(slide, text, left, top, width=Inches(1.5), color=ACCENT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(
        min(255, color[0] // 6 + 0x0a),
        min(255, color[1] // 6 + 0x0a),
        min(255, color[2] // 6 + 0x0a)
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, 12, True, color=color)
    return shape


def add_pyramid_level(slide, text, left, top, width, height=Inches(0.5), color=ACCENT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(
        min(255, color[0] // 6 + 0x0a),
        min(255, color[1] // 6 + 0x0a),
        min(255, color[2] // 6 + 0x0a)
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, 11, True, color=TEXT_MAIN)
    return shape


def add_simple_table(slide, headers, rows, left, top, width, col_widths=None):
    """Add a styled table to the slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.4 * num_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        set_run(run, h, 12, True, color=ACCENT_RED)
        p.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    # Data rows
    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            set_run(run, str(cell_text), 11, False, color=TEXT_SUB)
            p.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_BG if r_idx % 2 == 0 else TABLE_ROW_ALT

    # Remove table borders (set to thin dark lines)
    from pptx.oxml.ns import qn
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = tbl.makeelement(qn('a:tblPr'), {})
        tbl.insert(0, tblPr)
    tblPr.set('bandRow', '0')

    return table_shape


# ──────────────────────────────────────────────────────
# SLIDE BUILDERS
# ──────────────────────────────────────────────────────

def make_tech_slide(emoji, title_text_str, subtitle_str, subtitle_color,
                    part_num, left_title, left_body, right_tags, stats,
                    bridge_start, bridge_end_text, source_text, icon_color=ACCENT_RED):
    """Helper for the 8 tech analysis slides (slides 5-12)."""
    slide = add_slide()
    bridge_text(slide, bridge_start, Inches(0.15))
    part_marker(slide, f'PART 2 — 기술별 분석 ({part_num}/8)', Inches(0.65))

    # Icon
    icon_box(slide, emoji, SW/2 - Inches(0.6), Inches(1.1), Inches(1.2), icon_color)

    # Title
    tf = add_textbox(slide, Inches(0.5), Inches(2.4), SW - Inches(1), Inches(0.6))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, title_text_str, 30, True, color=TEXT_MAIN)

    # Subtitle
    tf2 = add_textbox(slide, Inches(0.5), Inches(3.0), SW - Inches(1), Inches(0.5))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    set_run(run2, subtitle_str, 18, True, color=subtitle_color)

    # Two-column content
    col_left = Inches(1)
    col_right = Inches(7)
    col_w = Inches(5.5)

    # Left column
    tf_l = add_textbox(slide, col_left, Inches(3.7), col_w, Inches(0.4))
    p = tf_l.paragraphs[0]
    run = p.add_run()
    set_run(run, left_title, 16, True, color=TEXT_DIM)

    tf_lb = add_textbox(slide, col_left, Inches(4.1), col_w, Inches(1.5))
    p = tf_lb.paragraphs[0]
    p.line_spacing = Pt(24)
    # handle **bold** in left body
    parts = left_body.split('**')
    for j, part in enumerate(parts):
        if not part:
            continue
        run = p.add_run()
        if j % 2 == 1:
            set_run(run, part, 14, True, color=TEXT_MAIN)
        else:
            set_run(run, part, 14, False, color=TEXT_SUB)

    # Right column - tags
    tf_r = add_textbox(slide, col_right, Inches(3.7), col_w, Inches(0.4))
    p = tf_r.paragraphs[0]
    run = p.add_run()
    set_run(run, '핵심 키워드', 16, True, color=TEXT_DIM)

    tag_colors = [ACCENT_RED, ACCENT_CYAN, ACCENT_GREEN, ACCENT_GOLD, ACCENT_PURPLE]
    tag_left = col_right
    for i, tag in enumerate(right_tags):
        c = tag_colors[i % len(tag_colors)]
        add_tag(slide, tag, tag_left, Inches(4.15), c)
        tag_left += Inches(1.5)

    # Stats
    for i, stat in enumerate(stats):
        tf_s = add_textbox(slide, col_right, Inches(4.65 + i * 0.4), col_w, Inches(0.4))
        p = tf_s.paragraphs[0]
        parts = stat.split('**')
        for j, part in enumerate(parts):
            if not part:
                continue
            run = p.add_run()
            if j % 2 == 1:
                set_run(run, part, 12, True, color=icon_color)
            else:
                set_run(run, part, 12, False, color=TEXT_DIM)

    bridge_end(slide, bridge_end_text, Inches(6.6))
    source_footer(slide, source_text, Inches(7.0))
    return slide


# ══════════════════════════════════════════
# BUILD ALL 62 SLIDES
# ══════════════════════════════════════════

# ── SLIDE 1: COVER ──
s = add_slide()
tf = add_textbox(s, Inches(0.5), Inches(1.8), SW - Inches(1), Inches(0.8))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '📖  ⚙️  🏭  💻  🌐  📱  🧠  🤖'
run.font.size = Pt(36)

tf2 = add_textbox(s, Inches(0.5), Inches(2.8), SW - Inches(1), Inches(1.5))
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
set_run(run2, '기술진보를 관통하는', 40, True, color=TEXT_MAIN)
p3 = tf2.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
set_run(run3, '하나의 키워드', 44, True, color=ACCENT_RED)

subtitle_text(s, '인쇄술에서 휴머노이드까지 — 문명 진화의 일관된 방향', Inches(4.5), 20, TEXT_DIM)
subtitle_text(s, '2025', Inches(5.3), 16, TEXT_FAINT)


# ── SLIDE 2: TOC ──
s = add_slide()
title_text(s, '목차', Inches(0.5), 32)
toc_items = [
    ('01', '도입 — 질문 던지기', 'p.3'),
    ('02', '8대 기술 혁신 분석', 'p.4–11'),
    ('03', '관통 키워드 도출', 'p.12–20'),
    ('04', '혁신의 드라이버', 'p.21–28'),
    ('05', '결론 — 최종 정리', 'p.29–30'),
    ('06', '일의 재정의 — 레버리지가 무한에 가까워지면', 'p.31–37'),
    ('07', '권력의 이동 — 기술이 바꾸는 지배의 구조', 'p.38–50'),
    ('08', '양극화 전망 — 기초 평준화와 메타 권력의 집중', 'p.51–62'),
]
for i, (num, label, pages) in enumerate(toc_items):
    y = Inches(1.6 + i * 0.65)
    tf = add_textbox(s, Inches(3), y, Inches(7), Inches(0.55))
    p = tf.paragraphs[0]
    run_num = p.add_run()
    set_run(run_num, num + '  ', 20, True, color=ACCENT_RED)
    run_label = p.add_run()
    set_run(run_label, label, 17, False, color=TEXT_SUB)
    run_pages = p.add_run()
    set_run(run_pages, '   ' + pages, 13, False, color=TEXT_FAINT)


# ── SLIDE 3: EXECUTIVE SUMMARY ──
s = add_slide()
part_marker(s, 'EXECUTIVE SUMMARY', Inches(0.3))
title_text(s, '이 강의가 말하고자 하는 것', Inches(0.7), 28)

exec_items = [
    ('1', '인쇄술에서 휴머노이드까지, 인류의 모든 기술진보는 인간의 기억·힘·절차·계산·연결·판단·행동을 점차 기계에게 위임해 온 대리화(Agency Transfer)의 역사다.', ACCENT_RED),
    ('2', '이 대리화를 밀어붙이는 힘은 희소성 제거 압력이며, 생존·비용·속도·복잡성·대리 욕구가 겹겹이 쌓여 기술혁신을 가속한다.', ACCENT_CYAN),
    ('3', '대리화가 극한으로 가면 일의 본질이 바뀐다 — "하는 것"에서 "시키는 것"을 넘어, "정하는 것"만이 인간에게 남는다.', ACCENT_GOLD),
    ('4', '그 과정에서 권력은 의미의 독점에서 시스템의 독점으로 이동하고, 기초 자원은 평준화되지만 메타 권한은 오히려 소수에게 집중된다.', ACCENT_PURPLE),
    ('5', 'AI 시대의 적자생존은 가장 강한 자가 아니라, 가장 빠르게 시스템을 이해하고 지휘하는 자의 몫이다. Claude Code를 무기로 들어라.', ACCENT_RED),
]
for i, (num, text, color) in enumerate(exec_items):
    y = Inches(1.5 + i * 1.1)
    # number
    tf = add_textbox(s, Inches(1.5), y, Inches(0.6), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, num, 22, True, color=color)
    # text
    tf2 = add_textbox(s, Inches(2.2), y, Inches(9), Inches(1.0))
    p2 = tf2.paragraphs[0]
    p2.line_spacing = Pt(22)
    run2 = p2.add_run()
    set_run(run2, text, 14, False, color=TEXT_SUB)


# ── SLIDE 4: QUESTION ──
s = add_slide()
bridge_text(s, '600년간 8번의 대전환. 각각 달라 보이지만, 하나의 공통 패턴이 숨어 있다.', Inches(0.15))
icon_box(s, '❓', SW/2 - Inches(0.6), Inches(1.0), Inches(1.2), ACCENT_RED)
title_text(s, '인쇄술부터 휴머노이드까지', Inches(2.4), 30)
add_quote_box(s, '이 모든 혁신을 꿰뚫는\n**단 하나의 키워드**는\n무엇인가?',
              Inches(3), Inches(3.2), Inches(7.333), Inches(1.8), ACCENT_RED, 22)

tags_data = [('📖 1440', ACCENT_RED), ('⚙️ 1760', ACCENT_CYAN), ('🏭 1913', ACCENT_GREEN),
             ('💻 1945', ACCENT_GOLD), ('🌐 1991', ACCENT_PURPLE), ('📱 2007', ACCENT_RED),
             ('🧠 2012', ACCENT_CYAN), ('🤖 2024', ACCENT_GREEN)]
tag_start = Inches(2.2)
for text, col in tags_data:
    add_tag(s, text, tag_start, Inches(5.4), col, Inches(1.2))
    tag_start += Inches(1.3)


# ── SLIDES 5-12: 8 TECH SLIDES ──
tech_slides_data = [
    ('📖', '인쇄술  c. 1440', '기억과 지식의 외부화', ACCENT_RED, '1',
     '무엇을 바꿨나', '손으로 베끼던 지식을 **대량 복제 가능한 정보**로 전환. 구텐베르크의 활판 인쇄기는 책 1권 복제 시간을 수개월에서 수일로 단축했다.',
     ['복제성', '표준화', '지식 민주화'],
     ['1500년까지 유럽에서 약 **2,000만 권**의 책이 인쇄됨', '유럽 문해율: 1450년 ~5% → 1600년 ~20%'],
     '첫 번째 대전환 — 인간의 가장 오래된 병목은 \'기억의 한계\'였다',
     '지식은 해방되었다. 그러나 인간의 팔과 다리는 여전히 한계 안에 있었다 →',
     '출처: Eisenstein (1979), Febvre & Martin (1958), Houston (2002)', ACCENT_RED),

    ('⚙️', '산업혁명  c. 1760–1840', '근육과 에너지의 외부화', ACCENT_CYAN, '2',
     '무엇을 바꿨나', '인간·동물의 힘을 **증기/기계 동력**으로 대체. 와트의 증기기관(1776)이 방직·운송·광업의 생산성을 수십 배 끌어올렸다.',
     ['동력화', '기계화', '생산성 폭증'],
     ['영국 1인당 GDP: 1760년 대비 1840년 **약 2배** 증가', '영국 면직물 생산량: 1760→1840년 약 100배 증가'],
     '인쇄술이 지식의 복제를 풀었다면, 다음 병목은 \'근육의 한계\'였다',
     '기계가 힘을 대신했다. 그러나 숙련공의 감각과 절차는 여전히 사람에게 묶여 있었다 →',
     '출처: Maddison Project Database (2020), Mokyr (1990)', ACCENT_CYAN),

    ('🏭', '생산 컨베이어  c. 1913', '숙련과 공정의 외부화', ACCENT_GREEN, '3',
     '무엇을 바꿨나', '장인이 처음부터 끝까지 만들던 방식을 **공정 분해 + 흐름 최적화**로 전환. 포드의 조립 라인은 모델T 생산 시간을 12시간에서 93분으로 줄였다.',
     ['분업', '표준화', '시스템화', '흐름 최적화'],
     ['모델T 가격: $850(1908) → **$260**(1925)', '포드 일당 $5 정책(1914) → 미국 중산층 형성의 기폭제'],
     '산업혁명이 힘을 기계화했다면, 다음 과제는 \'숙련과 공정\'의 기계화였다',
     '공정은 표준화되었다. 그러나 계산과 정보 처리는 여전히 인간의 머리에 의존했다 →',
     '출처: Ford Motor Company Archives, Hounshell (1984)', ACCENT_GREEN),

    ('💻', '컴퓨터  c. 1945–', '계산과 형식적 사고의 외부화', ACCENT_GOLD, '4',
     '무엇을 바꿨나', '인간의 계산·사무·규칙 기반 처리를 **프로그래밍 가능한 기계**가 수행. ENIAC(1945)에서 현대 프로세서까지, 연산 속도는 **수조 배** 향상되었다.',
     ['연산 자동화', '규칙 기반 처리', '정형 판단'],
     ['무어의 법칙: 트랜지스터 수 **18개월마다 2배** (1965–2020)', 'ENIAC(1945): 30톤 → 현대 스마트폰: 수백만 배 빠른 연산력'],
     '컨베이어가 공정을 표준화했다면, 다음 병목은 \'계산과 사무\'의 속도였다',
     '계산은 자동화되었다. 그러나 각 컴퓨터는 고립되어 있었다 — 연결이 없었다 →',
     '출처: Moore (1965), Computer History Museum', ACCENT_GOLD),

    ('🌐', '인터넷  c. 1991–', '연결과 소통의 외부화', ACCENT_PURPLE, '5',
     '무엇을 바꿨나', '지역에 묶여 있던 정보를 **전 세계 즉시 연결**로 전환. 개인·기관 사이의 정보 흐름이 시스템화되고 거리의 제약이 사라졌다.',
     ['네트워크화', '탈중심화', '실시간성', '연결성'],
     ['인터넷 사용자: 1995년 1,600만 → 2024년 **54억 명**', '글로벌 전자상거래: 1998년 $10B → 2024년 $6.3T'],
     '컴퓨터가 계산을 자동화했다면, 다음 병목은 \'정보의 고립\'이었다',
     '세상이 연결되었다. 그러나 인터넷은 책상 위에 묶여 있었다 →',
     '출처: ITU, Internet World Stats (2024), Statista', ACCENT_PURPLE),

    ('📱', '스마트폰  c. 2007–', '생활 운영의 외부화', ACCENT_RED, '6',
     '무엇을 바꿨나', '컴퓨터가 책상에서 떨어져 나와 **몸에 붙어 다니는 시대**. 일상적 판단, 일정, 위치, 검색, 선택이 실시간으로 디지털화되었다.',
     ['개인화', '이동성', '상시접속', '즉시성'],
     ['전 세계 스마트폰 사용자: **약 69억 대** (2024)', '글로벌 앱 경제 규모: 연간 약 $935B (2024)'],
     '인터넷이 세상을 연결했다면, 다음 병목은 \'이동성\'이었다',
     '모든 것이 손 안에 들어왔다. 그러나 사용자가 직접 판단하고 선택해야 했다 →',
     '출처: Statista, Pew Research (2024), data.ai', ACCENT_RED),

    ('🧠', '인공지능 (AI)  c. 2012–', '인지와 판단의 외부화', ACCENT_CYAN, '7',
     '무엇을 바꿨나', '인간 고유라고 여겼던 **패턴 인식, 언어 생성, 판단 보조, 창작**을 기계가 수행. 딥러닝(2012)과 대규모 언어모델(2020~)이 인지 노동의 자동화를 가속했다.',
     ['인지 자동화', '추론 보조', '생성', '자율화'],
     ['글로벌 AI 시장: 2024년 약 **$1,840억** → 2030년 $8,270억', 'ChatGPT: 출시 2개월 만에 1억 사용자 — 역대 최고속 앱'],
     '스마트폰이 모든 것을 손 안에 넣었다면, 다음 병목은 \'판단\'이었다',
     '사고가 기계화되었다. 그러나 AI는 디지털 세계 안에 갇혀 있었다 — 몸이 없었다 →',
     '출처: Grand View Research (2024), McKinsey (2023), Reuters', ACCENT_CYAN),

    ('🤖', '휴머노이드 로봇  c. 2024–', '행동과 노동의 외부화', ACCENT_GREEN, '8',
     '무엇을 바꿨나', 'AI가 **물리적 몸**을 얻음. 디지털 자동화가 현실 세계의 이동·조작·환경 적응으로 확장. "생각하는 기계"에서 **"행동하는 기계"**로의 전환.',
     ['물리적 실행', '환경 적응', '자율 수행', '현장 노동 대체'],
     ['휴머노이드 시장: 2035년 **$380억** 규모 전망', 'Tesla Optimus 목표 단가: $20,000~$25,000'],
     'AI가 사고를 기계화했다면, 다음 병목은 \'물리적 실행\'이었다',
     '기억부터 행동까지 — 인간의 모든 기능이 외부화 가능해졌다. 이 8단계를 관통하는 키워드는? →',
     '출처: Goldman Sachs (2024), Tesla AI Day (2024), Boston Dynamics', ACCENT_GREEN),
]

for data in tech_slides_data:
    make_tech_slide(*data)


# ── SLIDE 13: 8-STAGE TABLE ──
s = add_slide()
bridge_text(s, '8개 기술을 나란히 놓으면, 일관된 패턴이 드러난다', Inches(0.15))
part_marker(s, 'PART 3 — 관통 키워드 도출', Inches(0.6))
title_text(s, '문명 진화 8단계 통합 테이블', Inches(1.0), 28)

headers = ['#', '기술', '외부화된 인간 능력', '핵심 키워드']
rows = [
    ['1', '📖 인쇄술', '기억 · 지식 전파', '복제, 표준화'],
    ['2', '⚙️ 산업혁명', '근육 · 에너지', '동력화, 기계화'],
    ['3', '🏭 컨베이어', '숙련 · 공정 감각', '분업, 시스템화'],
    ['4', '💻 컴퓨터', '계산 · 형식적 사고', '연산, 자동처리'],
    ['5', '🌐 인터넷', '연결 · 소통', '네트워크, 탈중심'],
    ['6', '📱 스마트폰', '생활 운영 · 즉시 판단', '개인화, 상시접속'],
    ['7', '🧠 AI', '인지 · 판단 · 창작', '추론, 생성'],
    ['8', '🤖 휴머노이드', '행동 · 현장 노동', '자율수행, 대리'],
]
add_simple_table(s, headers, rows, Inches(2), Inches(1.8), Inches(9.3),
                 [Inches(0.6), Inches(2.2), Inches(3.5), Inches(3)])
source_footer(s, '출처: 자체 분석 종합')


# ── SLIDE 14: KEYWORD ① SCALING ──
s = add_slide()
part_marker(s, '관통 키워드 후보 ①', Inches(0.2))
# Keyword box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(0.8), Inches(3.3), Inches(0.8)) if False else None
tf = add_textbox(s, Inches(3.5), Inches(0.7), Inches(6.3), Inches(0.9))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '확장 (Scaling)', 36, True, color=ACCENT_RED)

subtitle_text(s, '인간 능력의 배율을 끊임없이 높여온 역사', Inches(1.6), 18, TEXT_DIM)

cards = [
    ('📖', '지식 확장', '1명의 저술 → 수만 명이 읽음', ACCENT_RED),
    ('⚙️', '힘 확장', '1마력 → 수천 마력', ACCENT_CYAN),
    ('🧠', '지능 확장', '1명의 판단 → 초당 수억 건 처리', ACCENT_GOLD),
    ('🤖', '행동 확장', '1명의 노동 → 24시간 무인 가동', ACCENT_GREEN),
]
for i, (ico, title, desc, col) in enumerate(cards):
    x = Inches(1.5 + i * 2.8)
    add_card(s, x, Inches(2.5), Inches(2.5), Inches(2.2), ico, title, desc, col)
source_footer(s, '출처: Brynjolfsson & McAfee (2014)')


# ── SLIDE 15: KEYWORD ② EXTERNALIZATION ──
s = add_slide()
part_marker(s, '관통 키워드 후보 ②', Inches(0.2))
tf = add_textbox(s, Inches(3.5), Inches(0.7), Inches(6.3), Inches(0.9))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '외부화 (Externalization)', 36, True, color=ACCENT_CYAN)
subtitle_text(s, '인간 내부 능력을 도구·시스템 바깥으로 분리', Inches(1.6), 18, TEXT_DIM)

# Timeline flow
flow_items = ['🧠 기억', '📖 인쇄물', '💪 근육', '⚙️ 엔진', '🤔 판단', '🧠 AI', '🏃 행동', '🤖 로봇']
x = Inches(0.8)
for i, item in enumerate(flow_items):
    add_flow_step(s, item, x, Inches(3.0), Inches(1.3), ACCENT_CYAN if i % 2 == 0 else ACCENT_PURPLE)
    if i < len(flow_items) - 1:
        add_flow_arrow(s, x + Inches(1.3), Inches(3.0))
    x += Inches(1.55)

tf = add_textbox(s, Inches(1), Inches(4.2), Inches(11), Inches(1.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(26)
run = p.add_run()
set_run(run, '가장 철학적으로 정확한 표현 — 인간의 기능이 몸과 뇌 밖으로 빠져나가\n기술 시스템 안에 구현되는 과정', 15, False, True, TEXT_DIM)
source_footer(s, '출처: Stiegler (1998), McLuhan (1964)')


# ── SLIDE 16: KEYWORD ③ AGENCY TRANSFER ──
s = add_slide()
part_marker(s, '관통 키워드 후보 ③', Inches(0.2))
tf = add_textbox(s, Inches(3.5), Inches(0.7), Inches(6.3), Inches(0.9))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '대리화 (Agency Transfer)', 36, True, color=ACCENT_GREEN)
subtitle_text(s, '기술의 진화 = 도구 → 기계 → 시스템 → 행위자', Inches(1.6), 18, TEXT_DIM)

flow_data = [('도구\nTool', ACCENT_RED), ('기계\nMachine', ACCENT_CYAN), ('시스템\nSystem', ACCENT_PURPLE),
             ('네트워크\nNetwork', ACCENT_GOLD), ('모델\nModel', ACCENT_GREEN), ('행위자\nAgent', ACCENT_RED)]
x = Inches(1.2)
for i, (txt, col) in enumerate(flow_data):
    add_flow_step(s, txt.split('\n')[0], x, Inches(3.0), Inches(1.5), col)
    tf_sub = add_textbox(s, x, Inches(3.5), Inches(1.5), Inches(0.4))
    p = tf_sub.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, txt.split('\n')[1] if '\n' in txt else '', 10, False, color=TEXT_FAINT)
    if i < len(flow_data) - 1:
        add_flow_arrow(s, x + Inches(1.5), Inches(3.0))
    x += Inches(1.9)

tf = add_textbox(s, Inches(1), Inches(4.5), Inches(11), Inches(0.8))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '인간이 하던 기능을 점점 더 기계에게 위임하는 역사', 16, True, color=TEXT_SUB)
source_footer(s, '출처: Russell & Norvig (2020)')


# ── SLIDE 17: KEYWORD ④ COST COLLAPSE ──
s = add_slide()
part_marker(s, '관통 키워드 후보 ④', Inches(0.2))
tf = add_textbox(s, Inches(3.5), Inches(0.7), Inches(6.3), Inches(0.9))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '비용 붕괴 (Cost Collapse)', 36, True, color=ACCENT_GOLD)
subtitle_text(s, '각 기술은 항상 어떤 비용을 무너뜨렸다', Inches(1.6), 18, TEXT_DIM)

headers = ['기술', '붕괴된 비용', '변화']
rows = [
    ['📖 인쇄술', '지식 복제 비용', '책 가격 ~80% 하락'],
    ['⚙️ 산업혁명', '생산 비용', '면직물 가격 ~90% 하락'],
    ['🏭 컨베이어', '조립 비용', '모델T 가격 ~70% 하락'],
    ['💻 컴퓨터', '연산 비용', 'MIPS당 비용 10년마다 ~90% 하락'],
    ['🌐 인터넷', '전송/탐색 비용', '정보 전달 비용 → 거의 0'],
    ['📱 스마트폰', '접근 비용', '컴퓨팅 접근 → 주머니 속'],
    ['🧠 AI', '인지/창작 비용', 'GPT-4 토큰 비용 2년간 ~95% 하락'],
    ['🤖 휴머노이드', '현장 노동 비용', '목표: 인건비 대비 ~50% 절감'],
]
add_simple_table(s, headers, rows, Inches(2.5), Inches(2.2), Inches(8.3),
                 [Inches(2), Inches(3), Inches(3.3)])
source_footer(s, '출처: Our World in Data, Nordhaus (2007), OpenAI API Pricing History')


# ── SLIDE 18: FINAL KEYWORD SELECTION ──
s = add_slide()
part_marker(s, '최종 키워드 선정', Inches(0.2))
title_text(s, '왜 "대리화"가 가장 강력한가?', Inches(0.8), 28)

# VS boxes
vs_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.0), Inches(4), Inches(2.5)) if False else None
shape_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.0), Inches(4), Inches(2))
shape_l.fill.solid()
shape_l.fill.fore_color.rgb = CARD_BG
shape_l.line.color.rgb = ACCENT_CYAN
shape_l.line.width = Pt(1)
tf = shape_l.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(14)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '🛠️'
run.font.size = Pt(36)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
set_run(run2, '"더 잘 도와줘"', 20, True, color=TEXT_MAIN)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
set_run(run3, '과거 기술의 목적', 14, False, color=TEXT_DIM)

# Arrow
tf_arr = add_textbox(s, Inches(6.2), Inches(2.5), Inches(1), Inches(1))
p = tf_arr.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '→', 36, True, color=ACCENT_RED)

shape_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(2.0), Inches(4), Inches(2))
shape_r.fill.solid()
shape_r.fill.fore_color.rgb = CARD_BG
shape_r.line.color.rgb = ACCENT_RED
shape_r.line.width = Pt(2)
tf = shape_r.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(14)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '🤖'
run.font.size = Pt(36)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
set_run(run2, '"내 대신 해줘"', 20, True, color=ACCENT_RED)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
set_run(run3, 'AI·휴머노이드 시대', 14, False, color=TEXT_DIM)

tf = add_textbox(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(24)
run = p.add_run()
set_run(run, '휴머노이드가 들어오면 기술의 목적이 보조를 넘어\n자율적 대리 수행으로 완전히 이동한다.\n"단순 도구"가 아니라 대리 수행자(Agent)에 가까워진다.', 15, False, color=TEXT_SUB)


# ── SLIDE 19: 6-STAGE PATTERN ──
s = add_slide()
bridge_text(s, '대리화로 가는 길에는 반복되는 공식이 있다 — 모든 기술이 거치는 6단계 패턴', Inches(0.15))
part_marker(s, '기술 진화의 공통 공식', Inches(0.6))
title_text(s, '기술진보의 6단계 반복 패턴', Inches(1.0), 28)

steps = [('① 복제', ACCENT_RED), ('② 표준화', ACCENT_CYAN), ('③ 대량화', ACCENT_PURPLE),
         ('④ 연결화', ACCENT_GOLD), ('⑤ 자동화', ACCENT_GREEN), ('⑥ 자율화', ACCENT_RED)]
x = Inches(1)
for i, (txt, col) in enumerate(steps):
    add_flow_step(s, txt, x, Inches(2.5), Inches(1.6), col)
    if i < len(steps) - 1:
        add_flow_arrow(s, x + Inches(1.6), Inches(2.5))
    x += Inches(2)

tf = add_textbox(s, Inches(1), Inches(3.5), Inches(11.3), Inches(2))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(24)
run = p.add_run()
set_run(run, '모든 시대의 핵심 기술은 이 6단계를 반복합니다.\n인쇄술은 1→3단계를, 인터넷은 4단계를, AI와 휴머노이드는 5→6단계를 완성합니다.', 15, False, color=TEXT_DIM)


# ── SLIDE 20: 8-STAGE DIAGRAM ──
s = add_slide()
bridge_text(s, '이 6단계 공식을 8개 기술에 적용하면, 문명 진화의 타임라인이 완성된다', Inches(0.15))
title_text(s, '문명 진화 8단계 다이어그램', Inches(0.7), 28)

timeline_items = [
    ('📖', '기억의\n외부화', '1440', ACCENT_RED),
    ('⚙️', '근육의\n외부화', '1760', ACCENT_CYAN),
    ('🏭', '공정의\n외부화', '1913', ACCENT_GREEN),
    ('💻', '계산의\n외부화', '1945', ACCENT_GOLD),
    ('🌐', '연결의\n외부화', '1991', ACCENT_PURPLE),
    ('📱', '생활의\n외부화', '2007', ACCENT_RED),
    ('🧠', '판단의\n외부화', '2012', ACCENT_CYAN),
    ('🤖', '행동의\n외부화', '2024', ACCENT_GREEN),
]
x = Inches(0.6)
for i, (emoji, label, year, col) in enumerate(timeline_items):
    # emoji
    tf = add_textbox(s, x, Inches(1.8), Inches(1.3), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = emoji
    run.font.size = Pt(28)
    # label
    tf2 = add_textbox(s, x, Inches(2.35), Inches(1.3), Inches(0.7))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    set_run(run2, label, 11, True, color=col)
    # year
    tf3 = add_textbox(s, x, Inches(3.1), Inches(1.3), Inches(0.4))
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    set_run(run3, year, 10, False, color=TEXT_FAINT)
    # arrow
    if i < len(timeline_items) - 1:
        add_flow_arrow(s, x + Inches(1.2), Inches(2.1))
    x += Inches(1.55)

add_quote_box(s, '인류 기술사의 일관된 방향:\n**인간 기능의 점진적 외부화와 대리화**',
              Inches(2.5), Inches(4.2), Inches(8.3), Inches(1.4), ACCENT_RED, 20)


# ── SLIDE 21: COMPRESSION ──
s = add_slide()
bridge_text(s, '8개 기술, 4개 렌즈, 6단계 공식 — 이 모든 분석을 한 문장으로 압축하면', Inches(0.15))
part_marker(s, '기술사 전체를 한 줄로 압축하면', Inches(0.6))

add_quote_box(s, '인류의 기술사는\n**기억 · 힘 · 절차 · 계산 · 연결 · 판단 · 행동**을\n점차 기계에게 위임해 가는\n**대리화의 역사**다.',
              Inches(2), Inches(1.5), Inches(9.3), Inches(2.2), ACCENT_RED, 24)

# Compare grid
labels = [('철학적 표현', '외부화', ACCENT_CYAN), ('경제학적 표현', '비용 붕괴', ACCENT_GOLD),
          ('공학적 표현', '자동화', ACCENT_GREEN), ('미래산업적 표현', '대리화', ACCENT_RED)]
for i, (lbl, val, col) in enumerate(labels):
    x = Inches(2 + i * 2.5)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), Inches(2.2), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = col
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, lbl, 11, False, color=TEXT_DIM)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    set_run(run2, val, 20, True, color=col)

bridge_end(s, '기술의 \'무엇\'은 정리했다. 이제 \'왜\' — 이 기술들을 밀어붙인 힘은 무엇인가? →', Inches(6.0))


# ── SLIDE 22: TOP-LEVEL DRIVER ──
s = add_slide()
part_marker(s, 'PART 4 — 혁신의 드라이버', Inches(0.2))
title_text(s, '최상위 드라이버', Inches(0.7), 24, TEXT_DIM)

tf = add_textbox(s, Inches(2.5), Inches(1.4), Inches(8.3), Inches(1))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '희소성 제거 압력', 40, True, color=ACCENT_RED)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
set_run(run2, 'Scarcity Removal Pressure', 18, False, color=TEXT_DIM)

tf2 = add_textbox(s, Inches(1.5), Inches(3.0), Inches(10.3), Inches(2))
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(28)
run = p.add_run()
set_run(run, '인간은 항상 희소한 것(시간, 노동, 에너지, 정보, 주의력, 판단력, 숙련, 이동성)을\n덜 희소하게 만들기 위해 기술을 발명한다.', 17, False, color=TEXT_SUB)

tf3 = add_textbox(s, Inches(1.5), Inches(4.8), Inches(10.3), Inches(1))
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(24)
run = p.add_run()
set_run(run, '"비싸고 느리고 부족한 인간 능력을,\n싸고 빠르고 풍부한 시스템 능력으로 바꾸려는 압력"', 15, False, True, TEXT_DIM)


# ── SLIDES 23-27: FIVE DRIVERS ──
driver_slides = [
    ('⚔️', '생존 압력', 'Survival Pressure', '가장 원초적인 혁신 동력',
     '희소성 제거 압력의 가장 원초적 형태 — 생존이 걸리면 혁신이 시작된다',
     [('🌾', '식량 생산', '인구 증가를 감당하기 위한 농업→산업 전환', ACCENT_RED),
      ('⚔️', '군사 경쟁', '전쟁이 기술을 밀어올림 (ENIAC, 인터넷)', ACCENT_CYAN),
      ('🏭', '산업 경쟁', '경쟁자를 이기기 위한 효율화 압력', ACCENT_GREEN),
      ('👴', '고령화 대응', '노동력 부족 → AI/로봇 대체 가속', ACCENT_GOLD)],
     '생존이 걸리면, 비용을 줄여야 한다 →',
     '출처: ARPANET 역사, UN Population Division', ACCENT_RED, '드라이버 ①'),

    ('💰', '비용 절감 압력', 'Cost Reduction', '"비용을 낮추는 자가 시장을 장악한다"',
     '생존 압력이 낳은 첫 번째 경제적 반응 — 더 싸게, 더 적은 자원으로',
     [('📉', '더 싸게', '원가의 지속적 하락', ACCENT_RED),
      ('⚡', '더 빨리', '생산·서비스 속도 향상', ACCENT_CYAN),
      ('👤', '더 적은 인력으로', '자동화로 인력 절감', ACCENT_GREEN),
      ('✅', '더 낮은 오류로', '정밀도·품질 향상', ACCENT_GOLD)],
     '비용을 줄이려면, 속도를 높여야 한다 →',
     '출처: Nordhaus (2007), AI Index Report (2024)', ACCENT_CYAN, '드라이버 ②'),

    ('⚡', '속도 경쟁', 'Speed Competition', '"속도 차이가 경쟁력 차이로 직결된다"',
     '비용 절감의 핵심 수단 — 같은 시간에 더 많이 처리하면 단가가 내려간다',
     [('💻', '계산 속도', 'ENIAC: 5,000연산/초 → GPU: 수조 연산/초', ACCENT_RED),
      ('🌐', '전달 속도', '편지: 수 주 → 메시지: 실시간', ACCENT_CYAN),
      ('📱', '반응 속도', '은행 방문: 수 시간 → 모바일: 수 초', ACCENT_GREEN),
      ('🧠', '판단 속도', '전문가 분석: 수 일 → AI: 수 초', ACCENT_GOLD)],
     '속도를 높이면, 복잡성이 폭발한다 →',
     '출처: Computer History Museum, Top500 (2024)', ACCENT_GREEN, '드라이버 ③'),

    ('🔀', '복잡성 대응', 'Complexity Management', '"복잡성이 인간의 처리 한계를 넘을 때, 기술이 대리한다"',
     '속도와 규모가 커지면 필연적으로 따라오는 것 — 인간의 처리 한계를 넘는 복잡성',
     [('📊', '데이터 폭증', '전 세계 데이터: 2010년 2ZB → 2025년 181ZB', ACCENT_RED),
      ('🔗', '공급망 복잡화', '글로벌 공급망 노드 수: 수만 개', ACCENT_CYAN),
      ('📜', '규제 증가', '금융 규제만 연간 수만 페이지 신규 발행', ACCENT_GREEN),
      ('⏱️', '실시간 의사결정', '밀리초 단위 거래, 초 단위 이상탐지', ACCENT_GOLD)],
     '복잡성이 인간 한계를 넘으면, 대리를 원하게 된다 →',
     '출처: IDC, Statista', ACCENT_PURPLE, '드라이버 ④'),

    ('🙋', '대리 욕구', 'Delegation Desire', '모든 드라이버가 수렴하는 지점',
     '모든 드라이버가 수렴하는 지점 — 인간의 가장 근본적인 욕구',
     [('😴', '"귀찮은 일"', '반복적이고 지루한 작업 위임', ACCENT_RED),
      ('⚠️', '"위험한 일"', '위험한 환경의 작업 대리', ACCENT_CYAN),
      ('🤯', '"어려운 일"', '복잡한 판단과 분석 위임', ACCENT_GREEN),
      ('🔄', '"반복적인 일"', '동일 패턴의 무한 반복', ACCENT_GOLD)],
     '이 5개 힘이 겹겹이 쌓여 기술혁신을 밀어붙인다 →',
     '출처: 자체 분석', ACCENT_GOLD, '드라이버 ⑤'),
]

for emoji, title, eng, sub, bridge_s, cards_data, bridge_e, source, accent, part_label in driver_slides:
    s = add_slide()
    bridge_text(s, bridge_s, Inches(0.15))
    part_marker(s, part_label, Inches(0.65))
    icon_box(s, emoji, SW/2 - Inches(0.6), Inches(1.1), Inches(1.2), accent)

    tf = add_textbox(s, Inches(0.5), Inches(2.4), SW - Inches(1), Inches(0.6))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, title, 30, True, color=accent)

    subtitle_text(s, sub, Inches(3.0), 16, TEXT_DIM)

    for i, (ico, card_title, card_desc, col) in enumerate(cards_data):
        x = Inches(1.5 + i * 2.8)
        add_card(s, x, Inches(3.7), Inches(2.5), Inches(2.0), ico, card_title, card_desc, col)

    bridge_end(s, bridge_e, Inches(6.3))
    source_footer(s, source)


# ── SLIDE 28: DRIVER PYRAMID ──
s = add_slide()
bridge_text(s, '5개 드라이버는 독립적이 아니라, 서로를 강화하는 계층 구조를 이룬다', Inches(0.15))
part_marker(s, '드라이버 계층 구조', Inches(0.6))
title_text(s, '혁신 드라이버 5-Level 피라미드', Inches(1.0), 28)

pyramid_data = [
    ('L1: 희소성 제거 (최상위 드라이버)', Inches(4), ACCENT_RED),
    ('L2: 경제·산업 — 비용 절감 · 생산성 · 규모의 경제 · 속도', Inches(5.5), ACCENT_CYAN),
    ('L3: 조직 — 표준화 · 자동화 · 리스크 감소 · 통제 가능성', Inches(7), ACCENT_GREEN),
    ('L4: 인간 심리 — 편의 추구 · 귀찮음 회피 · 위험 회피', Inches(8.5), ACCENT_GOLD),
    ('L5: 문명 — 인구 증가 · 도시화 · 전쟁 · 무역 확대 · 고령화', Inches(10), ACCENT_PURPLE),
]
y = Inches(2.0)
for text, width, color in pyramid_data:
    left = (SW - width) / 2
    add_pyramid_level(s, text, left, y, width, Inches(0.55), color)
    y += Inches(0.7)
source_footer(s, '출처: 자체 분석 — 기술혁신 드라이버 계층 분석')


# ── SLIDE 29: DRIVER EVOLUTION ──
s = add_slide()
bridge_text(s, '드라이버가 강해질수록, 기술의 역할도 진화한다 — 보조에서 자율로', Inches(0.15))
part_marker(s, '드라이버의 진화 방향', Inches(0.6))
title_text(s, '보조에서 자율 수행으로', Inches(1.0), 28)

steps = [('🛠️ 보조\nAssist', ACCENT_RED), ('⚙️ 자동화\nAutomate', ACCENT_CYAN),
         ('🤝 대리\nDelegate', ACCENT_GOLD), ('🤖 자율 수행\nAutonomous', ACCENT_GREEN)]
x = Inches(1.5)
for i, (txt, col) in enumerate(steps):
    lines = txt.split('\n')
    add_flow_step(s, lines[0], x, Inches(2.5), Inches(2.3), col)
    tf = add_textbox(s, x, Inches(3.05), Inches(2.3), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, lines[1], 12, False, color=TEXT_FAINT)
    if i < len(steps) - 1:
        add_flow_arrow(s, x + Inches(2.3), Inches(2.5))
    x += Inches(2.8)

# Two boxes
shape_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(4.2), Inches(3.5), Inches(1.2))
shape_l.fill.solid()
shape_l.fill.fore_color.rgb = CARD_BG
shape_l.line.color.rgb = ACCENT_CYAN
shape_l.line.width = Pt(1)
tf = shape_l.text_frame
tf.word_wrap = True
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '화이트칼라 병목 제거', 14, True, color=ACCENT_CYAN)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
set_run(run2, '🧠 AI', 18, True, color=TEXT_MAIN)

shape_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(4.2), Inches(3.5), Inches(1.2))
shape_r.fill.solid()
shape_r.fill.fore_color.rgb = CARD_BG
shape_r.line.color.rgb = ACCENT_GREEN
shape_r.line.width = Pt(1)
tf = shape_r.text_frame
tf.word_wrap = True
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '블루칼라/현장 병목 제거', 14, True, color=ACCENT_GREEN)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
set_run(run2, '🤖 휴머노이드', 18, True, color=TEXT_MAIN)

bridge_end(s, '기술의 \'무엇\'과 \'왜\'를 정리했다. 사회에 미치는 첫 번째 충격은 — \'일\'의 정의가 바뀌는 것이다 →', Inches(6.2))


# ── SLIDE 30: FINAL SUMMARY (PART 5) ──
s = add_slide()
bridge_text(s, 'Part 2~4를 종합하면, 기술진보는 관점에 따라 다르게 읽히지만 방향은 하나다', Inches(0.15))
part_marker(s, 'PART 5 — 결론', Inches(0.6))
title_text(s, '최종 정리', Inches(1.0), 28)

headers = ['관점', '한 단어', '의미']
rows = [
    ['🔬 철학', '외부화', '인간 기능을 몸·뇌 밖으로 분리'],
    ['💰 경제학', '비용 붕괴', '희소했던 능력의 가격을 0에 수렴시킴'],
    ['⚙️ 공학', '자동화', '반복·규칙·판단을 기계로 전환'],
    ['🚀 미래산업', '대리화', '인간 역할의 자율적 대리 수행'],
]
add_simple_table(s, headers, rows, Inches(2.5), Inches(1.8), Inches(8.3),
                 [Inches(1.5), Inches(2), Inches(4.8)])

add_quote_box(s, '기술혁신은\n**희소한 인간 능력**을\n**덜 희소하게** 만들려는 압력에서\n발생한다.',
              Inches(3), Inches(4.2), Inches(7.3), Inches(2), ACCENT_RED, 20)


# ── SLIDE 31: TRANSITION - INFINITE LEVERAGE ──
s = add_slide()
tf = add_textbox(s, Inches(0.5), Inches(1.2), SW - Inches(1), Inches(0.6))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '📖 → ⚙️ → 🏭 → 💻 → 🌐 → 📱 → 🧠 → 🤖'
run.font.size = Pt(24)
run.font.color.rgb = TEXT_DIM

title_text(s, '지금까지의 결론을 바탕으로', Inches(2.0), 24, TEXT_DIM)

add_quote_box(s, '레버리지가 **무한**에 가까워지면,\n**"일"**의 정의는\n어떻게 바뀌는가?',
              Inches(3), Inches(3.0), Inches(7.3), Inches(2), ACCENT_CYAN, 24)

tf = add_textbox(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '대리화가 극한으로 가면, 인간의 역할 자체가 재정의된다. 이제 답해보자 →', 14, False, True, TEXT_DIM)


# ── SLIDE 32: PART 6 START ──
s = add_slide()
part_marker(s, 'PART 6 — 일의 재정의', Inches(0.2))
icon_box(s, '⚖️', SW/2 - Inches(0.6), Inches(0.8), Inches(1.2), ACCENT_GOLD)
title_text(s, '레버리지 ∞ 시대의 전제', Inches(2.1), 28)
bridge_text(s, 'Part 2~4에서 확인한 대리화·비용붕괴·드라이버를 결합하면, 이런 전제가 도출된다', Inches(2.7))

cards = [
    ('📈', '산출물 한계 소멸', '1명이 만들어낼 수 있는 결과물의 상한이 사라진다', ACCENT_RED),
    ('🧠🤖', '이중 대리 완성', 'AI가 인지노동을, 휴머노이드가 물리노동을 모두 대리', ACCENT_CYAN),
    ('📉', '한계비용 → 0', '추가 1단위 생산·서비스의 비용이 0에 수렴', ACCENT_GREEN),
    ('🔄', '실행의 범용화', '"잘 하는 것"만으로는 차별화가 불가능해진다', ACCENT_GOLD),
]
for i, (ico, title, desc, col) in enumerate(cards):
    x = Inches(1.5 + i * 2.8)
    add_card(s, x, Inches(3.3), Inches(2.5), Inches(2.2), ico, title, desc, col)

tf = add_textbox(s, Inches(1), Inches(5.8), Inches(11.3), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '이 조건에서 — "일"이란 무엇이 되는가?', 16, True, color=ACCENT_RED)
source_footer(s, '출처: Rifkin (2014), Brynjolfsson & McAfee (2014)')


# ── SLIDE 33: HISTORICAL EVOLUTION OF WORK ──
s = add_slide()
part_marker(s, 'PART 6 — 일의 재정의', Inches(0.2))
bridge_text(s, '일의 정의는 과거에도 기술이 바뀔 때마다 재정의되었다 — 이번이 처음이 아니다', Inches(0.55))
title_text(s, '"일"의 정의는 이미 계속 바뀌어 왔다', Inches(1.0), 26)

headers = ['시대', '"일"의 정의', '가치 기준', '보상 방식']
rows = [
    ['🌾 농경시대', '생존을 위한 육체 활동', '수확량', '현물'],
    ['🏭 산업시대', '시간을 팔고 임금을 받는 것', '노동시간', '시급·월급'],
    ['💼 지식경제', '전문 지식·판단을 제공하는 것', '성과·전문성', '연봉·성과급'],
    ['🧠 AI시대', '???', '???', '???'],
]
add_simple_table(s, headers, rows, Inches(1.5), Inches(1.8), Inches(10.3),
                 [Inches(1.8), Inches(3.5), Inches(2.5), Inches(2.5)])

tf = add_textbox(s, Inches(1), Inches(4.5), Inches(11.3), Inches(2))
p = tf.paragraphs[0]
p.line_spacing = Pt(22)
run = p.add_run()
set_run(run, '주당 노동시간 변화: 농경 ~70h → 산업혁명 ~60h → 현대 ~40h → AI시대 ?', 13, False, color=TEXT_DIM)
p2 = tf.add_paragraph()
p2.line_spacing = Pt(22)
p2.space_before = Pt(8)
run2 = p2.add_run()
set_run(run2, '매 시대마다 기술이 이전 시대의 "일"을 대체했고, 인간은 한 단계 위의 역할로 이동했다.\nAI가 "지식·판단"까지 대리하는 시대에, 인간은 어디로 이동하는가?', 13, False, color=TEXT_DIM)
source_footer(s, '출처: Arendt (1958), Graeber (2018)')


# ── SLIDE 34: EXECUTION VALUE COLLAPSE ──
s = add_slide()
part_marker(s, 'PART 6 — 일의 재정의', Inches(0.2))
bridge_text(s, '이번에 달라지는 것: \'실행\' 자체의 경제적 가치가 0에 수렴한다', Inches(0.55))
title_text(s, '사라지는 것: "실행"의 가치 붕괴', Inches(1.0), 26)

cards = [
    ('💻', '코딩', 'AI가 수만 줄의 코드를 분 단위로 생성'),
    ('✍️', '글쓰기', '보고서, 기사, 마케팅 카피 자동 생성'),
    ('🎨', '디자인', '이미지, UI, 영상 — 프롬프트 한 줄'),
    ('📊', '분석', '데이터 분석, 패턴 인식 — AI가 빠르고 정확'),
    ('🏭', '조립·제조', '휴머노이드가 24시간 비정형 작업'),
    ('🚚', '물류·배송', '자율주행 + 로봇 — 무인 라스트마일'),
    ('🌐', '통번역', '실시간 음성·문서 번역, 전문 통역 수요 급감'),
    ('📞', 'AI 콜센터', '24시간 자연어 응대, 상담원 80% 대체 전망'),
]
for i, (ico, title, desc) in enumerate(cards):
    row = i // 4
    col = i % 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(1.8 + row * 1.8)
    add_card(s, x, y, Inches(2.8), Inches(1.5), ico, title, desc, ACCENT_RED)

add_quote_box(s, '"**잘 하는 것**"은 더 이상 경쟁 우위가 아니다.\n기계가 **더 잘** 하기 때문이다.',
              Inches(3), Inches(5.7), Inches(7.3), Inches(1.2), ACCENT_RED, 18)


# ── SLIDE 35: FOUR NEW DEFINITIONS OF WORK ──
s = add_slide()
part_marker(s, 'PART 6 — 일의 재정의', Inches(0.2))
bridge_text(s, '실행이 사라진 자리에 남는 것은 무엇인가? — 4가지 인간 고유의 역할', Inches(0.55))
title_text(s, '남는 것: "일"의 새로운 4가지 정의', Inches(1.0), 26)

defs = [
    ('🧭 1. 의미 부여 (Meaning Making)', '무엇을, 왜 할 것인가를 정하는 능력.\nAI는 "어떻게"에 강하지만, "왜"는 인간이 정한다.', ACCENT_RED),
    ('🔍 2. 문제 발견 (Problem Finding)', '풀어야 할 문제를 찾아내는 능력.\nAI는 주어진 문제를 풀지만, 문제 자체를 정의하는 건 인간이다.', ACCENT_CYAN),
    ('👁️ 3. 취향과 기준 (Taste & Standards)', 'AI의 산출물 중 무엇이 좋은지 판별하는 능력.\n1,000개의 결과물에서 "이것"을 고르는 안목.', ACCENT_GOLD),
    ('🤝 4. 책임과 신뢰 (Accountability)', '최종 결정과 책임을 지는 인간의 역할.\nAI는 추천하지만, 서명하고 책임지는 건 사람이다.', ACCENT_GREEN),
]
for i, (title, desc, col) in enumerate(defs):
    row = i // 2
    c = i % 2
    x = Inches(1 + c * 5.8)
    y = Inches(1.8 + row * 2.5)
    add_card_left_align(s, x, y, Inches(5.5), Inches(2.2), title, desc, col)

source_footer(s, '출처: Newport (2024), Kahneman (2011)')


# ── SLIDE 36: WORK FORMULA CHANGE ──
s = add_slide()
part_marker(s, 'PART 6 — 일의 재정의', Inches(0.2))
bridge_text(s, '이 4가지를 공식으로 바꾸면, 일의 방정식 자체가 변한다', Inches(0.55))
title_text(s, '일의 공식이 바뀐다', Inches(1.0), 28)

formulas = [
    ('🏭 과거 — 산업시대', '일 = 시간 × 숙련도 × 노력', ACCENT_CYAN, CARD_BG),
    ('💼 현재 — 지식경제', '일 = 지식 × 판단 × 실행력', ACCENT_GOLD, CARD_BG),
    ('🧠 미래 — AI시대', '일 = 방향 설정 × 기준 제시 × 책임 수용', ACCENT_RED, RGBColor(0x20, 0x0a, 0x12)),
]
for i, (label, formula, col, bg) in enumerate(formulas):
    y = Inches(1.8 + i * 1.4)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y, Inches(8.3), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = col
    shape.line.width = Pt(1.5 if i == 2 else 1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, label, 13, False, color=TEXT_DIM)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    set_run(run2, formula, 18, True, color=col)

add_quote_box(s, '"**무엇을 할 것인가**"에서\n"**무엇을 시킬 것인가**"를 넘어\n"**무엇을 정할 것인가**"로',
              Inches(3), Inches(5.7), Inches(7.3), Inches(1.5), ACCENT_GOLD, 18)


# ── SLIDE 37: VALUE SHIFT ──
s = add_slide()
part_marker(s, 'PART 6 — 일의 재정의', Inches(0.2))
bridge_text(s, '공식이 바뀌면, 가치의 위계도 바뀐다 — 어디에 서 있는가가 중요해진다', Inches(0.55))
title_text(s, '가치의 이동: "실행"에서 "방향"으로', Inches(1.0), 26)

pyramid_data = [
    ('🧭 방향 설정자 (Director) — 인간 고유 영역: "왜, 무엇을"', Inches(5), ACCENT_RED),
    ('🏗️ 설계자 (Architect) — AI 협업으로 증폭: "어떤 구조로"', Inches(6.5), ACCENT_CYAN),
    ('📋 관리자 (Manager) — 부분 자동화 가능: "잘 돌아가게"', Inches(8), ACCENT_GOLD),
    ('⚙️ 실행자 (Executor) — 자동화/대체 가능 ⚠️: "시킨 대로"', Inches(9.5), ACCENT_PURPLE),
]
y = Inches(1.8)
for text, width, color in pyramid_data:
    left = (SW - width) / 2
    add_pyramid_level(s, text, left, y, width, Inches(0.6), color)
    y += Inches(0.75)

# Two comparison boxes
shape_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.3), Inches(4), Inches(1))
shape_l.fill.solid()
shape_l.fill.fore_color.rgb = CARD_BG
shape_l.line.color.rgb = TEXT_FAINT
shape_l.line.width = Pt(1)
tf = shape_l.text_frame
tf.word_wrap = True
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '과거의 핵심 질문', 11, False, color=TEXT_DIM)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
set_run(run2, '"이것을 잘 할 수 있는가?"', 15, True, color=TEXT_SUB)

shape_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.3), Inches(4), Inches(1))
shape_r.fill.solid()
shape_r.fill.fore_color.rgb = CARD_BG
shape_r.line.color.rgb = ACCENT_RED
shape_r.line.width = Pt(1.5)
tf = shape_r.text_frame
tf.word_wrap = True
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '미래의 핵심 질문', 11, False, color=TEXT_DIM)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
set_run(run2, '"이것을 해야 하는가?"', 15, True, color=ACCENT_RED)

source_footer(s, '출처: Drucker (1999), Huang (2024), WEF (2023)')


# ── SLIDE 38: PART 6 FINAL ANSWER ──
s = add_slide()
part_marker(s, 'Part 6 최종 답변', Inches(0.5))

add_quote_box(s, '일은\n"하는 것"에서\n"시키는 것"을 넘어\n**"정하는 것"**으로 바뀐다.',
              Inches(2.5), Inches(1.5), Inches(8.3), Inches(2.5), ACCENT_RED, 26)

tf = add_textbox(s, Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.2))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(24)
run = p.add_run()
set_run(run, '실행은 기계가, 판단은 AI가 대리하는 시대에\n인간에게 남는 "일"은 방향을 정하고, 기준을 세우고, 책임을 지는 것이다.', 16, False, color=TEXT_SUB)

# Transition box
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(5.8), Inches(7.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x14, 0x0a, 0x1e)
shape.line.color.rgb = ACCENT_PURPLE
shape.line.width = Pt(1.5)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(22)
run = p.add_run()
set_run(run, '그런데 여기서 한 가지 더:\n기술이 바뀌면, "누가" 세상을 움직이는가?\n권력은 어디로 이동하는가? →', 14, False, color=ACCENT_PURPLE)


# ── SLIDE 39: PART 7 START - POWER SHIFT ──
s = add_slide()
part_marker(s, 'PART 7 — 권력의 이동', Inches(0.2))
icon_box(s, '👑', SW/2 - Inches(0.6), Inches(0.8), Inches(1.2), ACCENT_PURPLE)
title_text(s, '기술이 바뀌면, 누가 세상을 움직이는가?', Inches(2.1), 26)
bridge_text(s, '기술이 인간 기능을 대리하고, 일의 정의가 바뀌면, 권력 구조도 반드시 재편된다', Inches(2.7))

tf = add_textbox(s, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(24)
run = p.add_run()
set_run(run, '기술진보는 단순히 도구를 바꾸는 것이 아니다.\n사회의 희소한 조정수단(control medium)을 바꾸고,\n그것을 장악한 계층에게 권력을 이동시킨다.', 16, False, color=TEXT_SUB)

# Power flow
flow = ['⛪ 성직자', '👑 군주', '🏛️ 정치인', '🏦 자본가', '💻 테크노크라시', '🔮 ???']
x = Inches(0.8)
for i, txt in enumerate(flow):
    col = [ACCENT_RED, ACCENT_CYAN, ACCENT_GREEN, ACCENT_GOLD, ACCENT_PURPLE, ACCENT_RED][i]
    add_flow_step(s, txt, x, Inches(5.5), Inches(1.7), col)
    if i < len(flow) - 1:
        add_flow_arrow(s, x + Inches(1.7), Inches(5.5))
    x += Inches(2.05)


# ── SLIDE 40: COMPLETE POWER MAP ──
s = add_slide()
bridge_text(s, '6번의 권력 교체에는 공통 패턴이 있다: 핵심 조정수단이 바뀌면 권력자가 바뀐다', Inches(0.15))
part_marker(s, 'PART 7 — 권력의 이동', Inches(0.6))
title_text(s, '권력 이동 전체 지도', Inches(1.0), 26)

headers = ['시대', '핵심 기술', '희소 자원', '권력자', '정당성']
rows = [
    ['전근대', '문자·종교체계', '의미·해석권', '⛪ 성직자', '"신이 원한다"'],
    ['왕조/제국', '군사·행정', '폭력·징세', '👑 군주', '"내가 지킨다"'],
    ['근대 국민국가', '인쇄·선거', '여론·대표성', '🏛️ 정치인', '"국민이 선택했다"'],
    ['산업자본주의', '공장·금융', '자본·생산수단', '🏦 자본가', '"시장이 배분한다"'],
    ['디지털 네트워크', '인터넷·데이터', '플랫폼·알고리즘', '💻 테크노크라트', '"데이터가 말한다"'],
    ['AI-로봇 시대', 'AI·에이전트·로봇', '모델·인프라·규칙', '🔮 시스템 주권자', '"시스템이 최적화했다"'],
]
add_simple_table(s, headers, rows, Inches(0.8), Inches(1.8), Inches(11.7),
                 [Inches(1.8), Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.7)])

tf = add_textbox(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '권력은 "사람"이 아니라 "희소한 매개체"를 장악한 자에게 이동한다', 14, True, color=ACCENT_RED)
source_footer(s, '출처: Weber (1922), Castells (2010)')


# ── SLIDES 41-45: POWER SHIFT STAGES ──
power_slides = [
    ('⛪', '성직자 — 의미의 독점', '권력 이동 ①',
     '글을 읽는 사람이 극소수. 경전 해석권이 집중.\n"왜 살아야 하는가"의 답을 가진 자가 사회를 정당화했다.',
     ['문해력 독점', '우주론 해석', '죄와 구원의 정의'],
     '성직자는 "세계 해석 API"를 독점한 자',
     '📖 인쇄술이 이 독점을 붕괴시켰다 → 종교개혁, 공론장 형성',
     '출처: Weber (1905), Eisenstein (1979)', ACCENT_RED),

    ('👑', '군주 — 강제의 독점', '권력 이동 ②',
     '행정·군사·징세 체계 강화.\n권력의 중심이 신의 대리인에서 폭력의 중앙집중자로 이동.',
     ['상비군', '세금 징수', '영토 관리', '관료제'],
     '군주는 "물리적 실행 엔진"을 독점한 자',
     '⚙️ 산업혁명이 경제력 중심을 이동 → 자본의 힘이 왕권을 압도',
     '출처: Tilly (1992), Anderson (1974)', ACCENT_CYAN),

    ('🏛️', '정치인 — 정당성의 독점', '권력 이동 ③',
     '인쇄술·신문·교육·철도·대중 징병이 결합하여 국민국가 등장.\n여론을 조직하고 집단 정체성을 설계하는 능력이 권력의 핵심.',
     ['선거 제도', '정당 조직', '대중 연설', '법률 제정'],
     '정치인은 "집단 서사와 대중 동원 플랫폼"을 장악한 자',
     '🏭 산업 자본이 실질적 삶의 권력을 가져감 → 일자리·임금을 자본이 결정',
     '출처: Habermas (1962), Anderson (1983)', ACCENT_GREEN),

    ('🏦', '자본가 — 생산의 독점', '권력 이동 ④',
     '공장·철도·석탄·전기·대량생산·금융.\n이 모든 것은 선거가 아니라 자본 축적이 있어야 움직인다.',
     ['일자리 창출/파괴', '임금 결정', '도시 형성', '금융 선점'],
     '자본가는 "사회적 에너지 흐름(생산·투자·고용)"을 장악한 자',
     '🌐 인터넷이 중개·유통·정보의 독점을 재편 → 플랫폼이 자본보다 강력한 관문',
     '출처: Marx (1867), Piketty (2014)', ACCENT_GOLD),

    ('💻', '테크노크라시 — 흐름의 독점', '권력 이동 ⑤',
     '직접 명령하지 않아도 되는 권력.\n보이지 않는 규칙이 행동을 지배한다.',
     ['검색 순위', '추천 알고리즘', '클라우드/API', '신원/결제'],
     '테크노크라트는 "행동 가능성의 경계선"을 설계한다',
     '입법보다 빠르고 · 시장보다 은밀하며 · 군대보다 일상 깊숙이 침투',
     '출처: Zuboff (2019), Srnicek (2017)', ACCENT_PURPLE),
]

for emoji, title, part_label, desc, tags, insight, transition, source, color in power_slides:
    s = add_slide()
    part_marker(s, part_label, Inches(0.2))
    icon_box(s, emoji, SW/2 - Inches(0.6), Inches(0.7), Inches(1.2), color)
    title_text(s, title, Inches(2.0), 26, color)

    tf = add_textbox(s, Inches(1.5), Inches(2.8), Inches(5.5), Inches(2))
    p = tf.paragraphs[0]
    p.line_spacing = Pt(22)
    run = p.add_run()
    set_run(run, desc, 14, False, color=TEXT_SUB)

    tag_x = Inches(1.5)
    tag_colors = [ACCENT_RED, ACCENT_CYAN, ACCENT_GREEN, ACCENT_GOLD]
    for i, tag in enumerate(tags):
        add_tag(s, tag, tag_x, Inches(4.5), tag_colors[i % 4])
        tag_x += Inches(1.8)

    # Right column insight
    tf2 = add_textbox(s, Inches(7.5), Inches(2.8), Inches(5), Inches(1.5))
    p = tf2.paragraphs[0]
    p.line_spacing = Pt(22)
    run = p.add_run()
    set_run(run, insight, 15, True, color=TEXT_MAIN)

    tf3 = add_textbox(s, Inches(7.5), Inches(4.3), Inches(5), Inches(1))
    p = tf3.paragraphs[0]
    p.line_spacing = Pt(20)
    run = p.add_run()
    set_run(run, transition, 12, False, True, TEXT_DIM)

    source_footer(s, source)


# ── SLIDE 46: INVISIBILITY OF POWER ──
s = add_slide()
bridge_text(s, '이 5단계를 관통하는 메타 패턴: 권력은 점점 더 보이지 않게 된다', Inches(0.15))
part_marker(s, '핵심 통찰', Inches(0.6))
title_text(s, '권력의 비가시화', Inches(1.0), 28, ACCENT_RED)

tf = add_textbox(s, Inches(1), Inches(1.8), Inches(11.3), Inches(0.6))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '권력은 더 강해질수록 더 추상화되고, 더 인프라화되며, 더 보이지 않는다', 16, False, color=TEXT_SUB)

vis_items = [
    ('⛪ 제단·의식', '매우 가시적', ACCENT_GREEN),
    ('👑 왕관·군대', '가시적', ACCENT_CYAN),
    ('🏛️ 의회·선거', '비교적 가시적', ACCENT_GOLD),
    ('🏦 본사·금융', '불투명', ACCENT_RED),
    ('💻 코드·서버', '거의 비가시적', ACCENT_PURPLE),
]
x = Inches(1)
for i, (label, vis, col) in enumerate(vis_items):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.7), Inches(2.1), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = col
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, label, 13, True, color=TEXT_MAIN)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    set_run(run2, vis, 11, False, color=col)
    x += Inches(2.3)

add_quote_box(s, '**명령하는** 권력 → **설계하는** 권력 → **기본값(default)을 정하는** 권력',
              Inches(1.5), Inches(4.8), Inches(10.3), Inches(1), ACCENT_RED, 18)
source_footer(s, '출처: Lessig (2006) — "Code is Law"')


# ── SLIDE 47: LEGITIMACY EVOLUTION ──
s = add_slide()
bridge_text(s, '보이지 않는 권력은, 정당성의 언어도 바꾼다 — 도덕에서 성능으로', Inches(0.15))
part_marker(s, 'PART 7 — 권력의 이동', Inches(0.6))
title_text(s, '정당성의 진화', Inches(1.0), 28)

headers = ['권력자', '정당성의 언어', '본질']
rows = [
    ['⛪ 성직자', '"신이 원한다"', '도덕'],
    ['👑 군주', '"내가 지킨다"', '폭력'],
    ['🏛️ 정치인', '"국민이 선택했다"', '대표성'],
    ['🏦 자본가', '"시장이 효율적으로 배분한다"', '효율'],
    ['💻 테크노크라트', '"데이터가 그렇게 말한다"', '최적화'],
    ['🔮 시스템 주권자', '"시스템이 최적화했다"', '???'],
]
add_simple_table(s, headers, rows, Inches(2.5), Inches(1.7), Inches(8.3),
                 [Inches(2.5), Inches(3.5), Inches(2.3)])

add_quote_box(s, '정당성이 **가치**에서 **성능**으로 이동한다\n"최적화"는 겉보기엔 중립적이지만, 실제로는\n**누구의 목적함수(objective function)**인지가 숨겨져 있다',
              Inches(2), Inches(5.2), Inches(9.3), Inches(1.5), ACCENT_GOLD, 16)
source_footer(s, '출처: Morozov (2013), Scott (1998)')


# ── SLIDE 48: SYSTEM SOVEREIGNTY ──
s = add_slide()
bridge_text(s, '테크노크라시를 넘어 — AI·로봇 시대의 권력은 시스템 스택 전체의 장악에 달린다', Inches(0.15))
part_marker(s, '다음 단계', Inches(0.6))
title_text(s, '시스템 주권 (System Sovereignty)', Inches(1.0), 28)

tf = add_textbox(s, Inches(1), Inches(1.7), Inches(11.3), Inches(0.6))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, 'AI·에이전트·휴머노이드 시대의 권력은 "누가 시스템 스택을 통합 장악하는가"에 달린다', 15, False, color=TEXT_SUB)

pyramid_data = [
    ('🔮 안전 규칙 · 윤리 프레임 — 무엇이 허용되는가', Inches(5), ACCENT_RED),
    ('🧠 AI 모델 · 에이전트 프레임워크 — 무엇을 판단하고 실행하는가', Inches(6.5), ACCENT_CYAN),
    ('⚡ 컴퓨팅 · 데이터 · 통신 — GPU/클라우드 · 학습 데이터', Inches(8), ACCENT_GREEN),
    ('🔐 인증/ID · 결제 레일 — 누가 참여하고 거래하는가', Inches(9.5), ACCENT_GOLD),
    ('🤖 로봇 OS · 센서 · 물리 인프라 — 현실 세계에서 무엇이 실행되는가', Inches(11), ACCENT_PURPLE),
]
y = Inches(2.5)
for text, width, color in pyramid_data:
    left = (SW - width) / 2
    add_pyramid_level(s, text, left, y, width, Inches(0.55), color)
    y += Inches(0.7)

tf = add_textbox(s, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '이 전체 스택을 묶는 주체가 다음 시대의 실질 권력자가 된다', 14, True, color=ACCENT_RED)
source_footer(s, '출처: 자체 분석; 빅테크 5사 AI 투자 총액 2024년 약 $200B+')


# ── SLIDE 49: FIVE FLOWS ──
s = add_slide()
bridge_text(s, '권력 이동 분석에서 도출되는 5가지 구조적 흐름', Inches(0.15))
part_marker(s, 'PART 7 — 정리', Inches(0.6))
title_text(s, '권력 이동에서 보이는 5가지 흐름', Inches(1.0), 26)

flows = [
    '① 권력은 "사람"이 아니라 "희소한 매개체"를 따라 이동한다',
    '② 기술은 기존 권력의 희소성 기반을 무너뜨린다',
    '③ "보이는 통치"에서 "보이지 않는 운영"으로 이동한다',
    '④ 정당성이 "가치"에서 "성능"으로 이동한다',
    '⑤ 미래 권력 = 인프라 + 모델 + 규칙의 통합 장악',
]
colors = [ACCENT_RED, ACCENT_CYAN, ACCENT_GREEN, ACCENT_GOLD, ACCENT_PURPLE]
for i, (text, col) in enumerate(zip(flows, colors)):
    y = Inches(1.8 + i * 1.0)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(10.3), Inches(0.75))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = col
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, text, 15, True, color=TEXT_MAIN)


# ── SLIDE 50: WARNING ──
s = add_slide()
bridge_text(s, '이 흐름이 계속되면, 가장 위험한 결과: 민주주의의 작동 범위 밖으로 권력이 이동', Inches(0.15))
part_marker(s, '경고', Inches(0.6))
title_text(s, '권력이 민주주의 밖으로 이동할 위험', Inches(1.0), 26, ACCENT_RED)

cards = [
    ('🗳️', '선출되지 않았다', '유권자가 뽑지 않은 자들이 설계'),
    ('⌨️', '코드로 작동한다', '법률이 아닌 알고리즘이 규칙'),
    ('🌍', '국경을 넘는다', '국가 단위 규제를 초월'),
    ('⚡', '속도가 빠르다', '입법 속도를 압도'),
    ('❓', '책임이 불명확하다', '"알고리즘이 결정했다"'),
    ('😊', '자발적 복종', '편의 때문에 기꺼이 의존'),
]
for i, (ico, title, desc) in enumerate(cards):
    row = i // 3
    col = i % 3
    x = Inches(1.5 + col * 3.6)
    y = Inches(1.8 + row * 2.0)
    add_card(s, x, y, Inches(3.2), Inches(1.7), ico, title, desc, ACCENT_RED)

add_quote_box(s, '다음 시대의 지배는\n강압보다\n**"의존성"**을 통해 이뤄진다',
              Inches(3), Inches(5.8), Inches(7.3), Inches(1.2), ACCENT_RED, 18)


# ── SLIDE 51: MIDDLE SUMMARY ──
s = add_slide()
part_marker(s, '지금까지의 정리', Inches(0.5))

summary_items = [
    ('🔧 기술사의 본질', '인간 기능의 대리화 역사', ACCENT_RED),
    ('💼 일의 재정의', '"하는 것"에서 "정하는 것"으로', ACCENT_CYAN),
    ('👑 권력의 이동', '"의미의 독점"에서 "시스템의 독점"으로', ACCENT_PURPLE),
]
for i, (label, desc, col) in enumerate(summary_items):
    y = Inches(1.2 + i * 1.5)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y, Inches(8.3), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = col
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, label, 16, True, color=col)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    set_run(run2, desc, 18, True, color=TEXT_MAIN)

# Transition question
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.7), Inches(8.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1a, 0x18, 0x0a)
shape.line.color.rgb = ACCENT_GOLD
shape.line.width = Pt(1.5)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(12)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '이 세 가지를 종합하면\n양극화는 어떻게 되는가? →', 18, True, color=ACCENT_GOLD)


# ── SLIDE 52: PART 8 START ──
s = add_slide()
part_marker(s, 'PART 8 — 양극화 전망', Inches(0.2))
icon_box(s, '⚖️', SW/2 - Inches(0.6), Inches(0.7), Inches(1.2), ACCENT_GOLD)
title_text(s, '기술 + 일 + 권력을 종합하면 양극화는?', Inches(2.0), 24)

add_quote_box(s, '미래의 양극화는\n"부자 vs 가난한 자"가 아니라\n**"시스템 지배력을 가진 소수"**와\n**"시스템에 의해 평가·대체되는 다수"** 사이의\n구조적 격차로 재편된다',
              Inches(2), Inches(3.0), Inches(9.3), Inches(2.5), ACCENT_RED, 18)

tf = add_textbox(s, Inches(1), Inches(5.8), Inches(11.3), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '그러나 — 이것은 단순한 비관론이 아니다. 표면과 심층을 구분해야 한다.', 14, False, True, TEXT_DIM)


# ── SLIDE 53: WHY DIFFERENT ──
s = add_slide()
part_marker(s, 'PART 8 — 양극화 전망', Inches(0.2))
bridge_text(s, '양극화는 늘 있었다. 그러나 이번에는 \'대체 범위\'가 다르다', Inches(0.55))
title_text(s, '왜 이번 양극화는 과거와 다른가?', Inches(1.0), 26)

# Left box - past
shape_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.0), Inches(5), Inches(3.5))
shape_l.fill.solid()
shape_l.fill.fore_color.rgb = CARD_BG
shape_l.line.color.rgb = ACCENT_CYAN
shape_l.line.width = Pt(1)
tf = shape_l.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(12)
p = tf.paragraphs[0]
run = p.add_run()
set_run(run, '과거: 일부 계층만 대체', 16, True, color=ACCENT_CYAN)
for line in ['', '⚙️ 산업혁명 → 육체노동 대체', '💻 디지털혁명 → 중간 사무직 대체']:
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run2 = p2.add_run()
    set_run(run2, line, 14, False, color=TEXT_SUB)

# Right box - now
shape_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3.5))
shape_r.fill.solid()
shape_r.fill.fore_color.rgb = RGBColor(0x20, 0x0a, 0x12)
shape_r.line.color.rgb = ACCENT_RED
shape_r.line.width = Pt(2)
tf = shape_r.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(12)
p = tf.paragraphs[0]
run = p.add_run()
set_run(run, '이번: 동시 다발 재편', 16, True, color=ACCENT_RED)
lines = ['', '🧠 AI → 인지노동 (사무·분석·창작·상담)', '🤖 휴머노이드 → 물리노동 (제조·물류·현장)',
         '', '= 이중 대리 완성', '', '화이트칼라 + 블루칼라 + 중간관리 + 전문직이 동시에 재편']
for line in lines:
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    if '이중 대리' in line:
        set_run(run2, line, 14, True, color=ACCENT_RED)
    else:
        set_run(run2, line, 13, False, color=TEXT_SUB)

source_footer(s, '출처: Autor (2024), Acemoglu & Restrepo (2020), IMF (2024)')


# ── SLIDE 54: 5-LAYER STRUCTURE ──
s = add_slide()
bridge_text(s, '이중 대리가 완성되면, 사회는 \'시스템과의 관계\'를 기준으로 재계층화된다', Inches(0.15))
part_marker(s, 'PART 8 — 양극화 전망', Inches(0.6))
title_text(s, '미래의 5계층 구조', Inches(1.0), 28)

pyramid_data = [
    ('① 시스템 소유층 — 모델·인프라·플랫폼 소유, 보이지 않는 준주권', Inches(5), ACCENT_RED),
    ('② 시스템 설계층 — AI연구·아키텍트·핵심엔지니어, 준자본가적 보상', Inches(6.5), ACCENT_CYAN),
    ('③ 시스템 지휘층 — 오케스트레이터·전략가·책임자, 기계를 지휘', Inches(8), ACCENT_GREEN),
    ('④ 시스템 보조층 ⚠️ — AI도구 활용·검수·예외처리, 가장 넓고 불안', Inches(9.5), ACCENT_GOLD),
    ('⑤ 시스템 종속층 — 대체·플랫폼노동화·이전소득 의존, 협상력 상실', Inches(11), ACCENT_PURPLE),
]
y = Inches(1.8)
for text, width, color in pyramid_data:
    left = (SW - width) / 2
    add_pyramid_level(s, text, left, y, width, Inches(0.6), color)
    y += Inches(0.75)

source_footer(s, '출처: 자체 분석 — AI시대 사회계층 모델; Frey & Osborne (2017)')


# ── SLIDE 55: BARGAINING POWER COLLAPSE ──
s = add_slide()
bridge_text(s, '계층화의 핵심 메커니즘은 소득이 아니라 \'협상력\'의 구조적 변화다', Inches(0.15))
part_marker(s, 'PART 8 — 양극화 전망', Inches(0.6))
title_text(s, '진짜 문제: 협상력 붕괴', Inches(1.0), 28, ACCENT_RED)

tf = add_textbox(s, Inches(1), Inches(1.7), Inches(11.3), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '양극화의 핵심은 소득이 아니라 존재의 bargaining power', 16, True, color=TEXT_SUB)

# VS
shape_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.5), Inches(4.5), Inches(2))
shape_l.fill.solid()
shape_l.fill.fore_color.rgb = CARD_BG
shape_l.line.color.rgb = ACCENT_CYAN
shape_l.line.width = Pt(1)
tf = shape_l.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(12)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '👤 사람 1명 채용', 18, True, color=TEXT_MAIN)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)
run2 = p2.add_run()
set_run(run2, '월급 + 보험 + 교육 + 관리\n+ 실수 리스크', 14, False, color=TEXT_DIM)

tf_arr = add_textbox(s, Inches(6.2), Inches(3.0), Inches(1), Inches(1))
p = tf_arr.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, 'vs', 20, True, color=ACCENT_RED)

shape_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(2.5), Inches(4.5), Inches(2))
shape_r.fill.solid()
shape_r.fill.fore_color.rgb = RGBColor(0x20, 0x0a, 0x12)
shape_r.line.color.rgb = ACCENT_RED
shape_r.line.width = Pt(2)
tf = shape_r.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(12)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '🤖 AI 1개 구독', 18, True, color=ACCENT_RED)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)
run2 = p2.add_run()
set_run(run2, '월 $20~200 + 24시간\n+ 무한 확장 + 무실수', 14, False, color=TEXT_DIM)

add_quote_box(s, '인간은 더 이상\n"필수 자원"이 아니라\n**"대체 가능한 옵션"**이 된다',
              Inches(3), Inches(5.0), Inches(7.3), Inches(1.3), ACCENT_RED, 18)

source_footer(s, '출처: Susskind (2020), Graeber (2018)')


# ── SLIDE 56: SURFACE EQUALIZATION ──
s = add_slide()
part_marker(s, '그러나 — 반론', Inches(0.2))
bridge_text(s, '그러나 모든 축이 양극화되는 것은 아니다 — 기술은 원래 희소성을 무너뜨린다', Inches(0.55))
title_text(s, '표면적 희소성은 오히려 약화된다', Inches(1.0), 26, ACCENT_GREEN)

cards = [
    ('💰', '자산: 도구의 민주화', '1인 기업, 1인 창작 스튜디오 가능.\n창업비용 하락, 소액 자본의 레버리지 확대.', ACCENT_GREEN),
    ('🧠', '역량: 기초 능력 평준화', '코딩, 글쓰기, 분석, 번역, 디자인 초안\n— AI가 중간 수준까지 급속 평준화.', ACCENT_GREEN),
    ('🔗', '접속권: 기본 접속 대중화', '기본 AI 모델·에이전트·자동화 도구도\n준공공재로 수렴 가능.', ACCENT_GREEN),
    ('🎭', '정체성: 다양화', '다중역할, 프로젝트형 삶, 취향 공동체,\n1인 다직업 — 다양화·파편화·유동화.', ACCENT_GREEN),
]
for i, (ico, title, desc, col) in enumerate(cards):
    row = i // 2
    c = i % 2
    x = Inches(1 + c * 5.8)
    y = Inches(1.8 + row * 2.5)
    add_card_left_align(s, x, y, Inches(5.5), Inches(2.2), f'{ico} {title}', desc, col)

source_footer(s, '출처: Diamandis & Kotler (2012), Rifkin (2014)')


# ── SLIDE 57: DEEP SCARCITY ──
s = add_slide()
part_marker(s, '그러나 동시에', Inches(0.2))
bridge_text(s, '표면은 평준화되지만, 그 위에서 새로운 \'메타 계층\'이 형성된다', Inches(0.55))
title_text(s, '심층 희소성은 재구성된다', Inches(1.0), 26, ACCENT_RED)

headers = ['축', '평준화되는 것', '집중되는 것']
rows = [
    ['💰 자산', '개인 도구, 1인 기업, 소규모 자동화', 'GPU 인프라, 데이터센터, 플랫폼, 결제레일'],
    ['🧠 역량', '실행 역량, 기초 지식, 초안 생산', '문제 정의, 목표 설계, 시스템 오케스트레이션'],
    ['🔗 접속권', '기본 모델, 범용 AI, 개인 비서', 'frontier 모델, 독점 데이터, 고신뢰 보안 권한'],
    ['🎭 정체성', '형태의 다양화, 다중 역할', '영향력·인정·지속성의 위계 재형성'],
]
add_simple_table(s, headers, rows, Inches(1.5), Inches(1.8), Inches(10.3),
                 [Inches(1.3), Inches(4), Inches(5)])

add_quote_box(s, '**"사용의 평등화"** + **"소유·설계·권한의 집중"**\n이 두 가지가 **동시에** 진행된다',
              Inches(2.5), Inches(5.2), Inches(8.3), Inches(1.2), ACCENT_RED, 18)


# ── SLIDE 58: FOUR META-AUTHORITIES ──
s = add_slide()
part_marker(s, 'PART 8 — 양극화 전망', Inches(0.2))
bridge_text(s, '평준화된 세계에서 진짜 희소해지는 것은 \'자원\'이 아니라 \'권한\'이다', Inches(0.55))
title_text(s, '진짜 희소해지는 것: 메타 권한', Inches(1.0), 26)

metas = [
    ('🎯 1. 목표 설정권', '누가 목적함수를 정하는가?\n무엇을 최적화할지, 누구의 편의를 기준으로 설계할지.', ACCENT_RED),
    ('🏗️ 2. 시스템 설계권', '누가 규칙과 구조를 만드는가?\n플랫폼 룰, API 구조, 알고리즘, 데이터 흐름.', ACCENT_CYAN),
    ('🔓 3. 예외 승인권', '자동화 사회에서 예외를 다루는 권력.\n승인/거절, 위험 해제, 수동 개입, 특별 대우.', ACCENT_GOLD),
    ('✍️ 4. 책임 인수권', '문제 발생 시 누가 책임지는가.\nAI가 오판하면? 로봇이 사고내면? 서명하는 건 사람.', ACCENT_GREEN),
]
for i, (title, desc, col) in enumerate(metas):
    row = i // 2
    c = i % 2
    x = Inches(1 + c * 5.8)
    y = Inches(1.8 + row * 2.5)
    add_card_left_align(s, x, y, Inches(5.5), Inches(2.2), title, desc, col)

tf = add_textbox(s, Inches(1), Inches(6.8), Inches(11.3), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '기술은 기초 자원은 풍부하게 만들지만, 이 4가지 메타 권한은 더 희소하게 만든다', 14, True, color=ACCENT_RED)


# ── SLIDE 59: SHORT-MEDIUM TERM ──
s = add_slide()
bridge_text(s, '이 구조가 시간축에서 어떻게 전개되는가 — 단기 전망', Inches(0.15))
part_marker(s, '시간축 전망 ①', Inches(0.6))
title_text(s, '단기~중기 2026–2035', Inches(1.0), 28)

cards = [
    ('📈', '생산성 급증', 'AI 도입 기업 생산성 20~40% 향상'),
    ('📉', '고용 구조 불안정', '중간 사무직·일부 전문직 직무 재편'),
    ('💸', '임금 압박', '성과 계량화 심화, 협상력 약화'),
    ('🏢', '조직 압축', '팀 규모 축소, 1인 다역할'),
    ('🎓', '스킬 전환 혼란', '"무엇을 배워야 하는가" 불확실'),
    ('📊', '상위층 초과보상', '기술+자본 결합층에 수익 집중'),
]
for i, (ico, title, desc) in enumerate(cards):
    row = i // 3
    col = i % 3
    x = Inches(1 + col * 3.8)
    y = Inches(1.8 + row * 1.8)
    add_card(s, x, y, Inches(3.5), Inches(1.5), ico, title, desc, ACCENT_CYAN)

add_quote_box(s, '"고용은 유지되는 것처럼 보이지만,\n**협상력과 임금의 질**이 무너지는 시기"',
              Inches(3), Inches(5.7), Inches(7.3), Inches(1.2), ACCENT_CYAN, 16)
source_footer(s, '출처: McKinsey Global Institute (2023), Goldman Sachs (2024)')


# ── SLIDE 60: MEDIUM-LONG TERM ──
s = add_slide()
bridge_text(s, '단기 충격이 지나면, 휴머노이드 본격화로 이중 충격이 합류한다', Inches(0.15))
part_marker(s, '시간축 전망 ②', Inches(0.6))
title_text(s, '중장기 2035–2050', Inches(1.0), 28)

tf = add_textbox(s, Inches(1), Inches(1.7), Inches(11.3), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_run(run, '휴머노이드 본격화 → 이중 충격의 합류', 18, True, color=TEXT_MAIN)

# Two boxes
shape_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.5), Inches(5), Inches(2.5))
shape_l.fill.solid()
shape_l.fill.fore_color.rgb = CARD_BG
shape_l.line.color.rgb = ACCENT_CYAN
shape_l.line.width = Pt(1.5)
tf = shape_l.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(12)
p = tf.paragraphs[0]
run = p.add_run()
set_run(run, '🧠 화이트칼라 충격 (AI)', 16, True, color=ACCENT_CYAN)
p2 = tf.add_paragraph()
p2.space_before = Pt(8)
run2 = p2.add_run()
set_run(run2, '사무 · 분석 · 상담 · 창작 · 관리\n→ 에이전트가 대부분 처리', 14, False, color=TEXT_SUB)

shape_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.5), Inches(5), Inches(2.5))
shape_r.fill.solid()
shape_r.fill.fore_color.rgb = CARD_BG
shape_r.line.color.rgb = ACCENT_GREEN
shape_r.line.width = Pt(1.5)
tf = shape_r.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(12)
p = tf.paragraphs[0]
run = p.add_run()
set_run(run, '🤖 블루칼라 충격 (휴머노이드)', 16, True, color=ACCENT_GREEN)
p2 = tf.add_paragraph()
p2.space_before = Pt(8)
run2 = p2.add_run()
set_run(run2, '물류 · 제조 · 점포 · 건설 · 시설관리\n→ 로봇이 현장 작업 수행', 14, False, color=TEXT_SUB)

# Alert box
shape_a = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.3), Inches(9.3), Inches(1.3))
shape_a.fill.solid()
shape_a.fill.fore_color.rgb = RGBColor(0x20, 0x0a, 0x12)
shape_a.line.color.rgb = ACCENT_RED
shape_a.line.width = Pt(2)
tf = shape_a.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(22)
run = p.add_run()
set_run(run, '두 충격이 합쳐지는 순간 — 단순 재교육만으로는 해결 어려움\n"배운다 해도 배운 것 자체가 자동화될 수 있다"', 14, True, color=ACCENT_RED)

source_footer(s, '출처: Goldman Sachs (2024), WEF (2023), 한국 통계청')


# ── SLIDE 61: FIVE VARIABLES ──
s = add_slide()
bridge_text(s, '양극화는 운명이 아니다 — 이 5개 변수의 방향에 따라 결과가 완전히 달라진다', Inches(0.15))
part_marker(s, 'PART 8 — 양극화 전망', Inches(0.6))
title_text(s, '양극화를 결정하는 5개 변수', Inches(1.0), 26)

variables = [
    ('① AI/로봇 소유 구조', '소수 기업 독점? 국가/공공 인프라화? 오픈 생태계? → 생산성 증가분의 배분을 결정', ACCENT_RED),
    ('② 노동의 재정의 속도', '직무 재설계가 가능한가? 인간 역할을 검수·책임·공감·현장으로 재배치할 수 있는가?', ACCENT_CYAN),
    ('③ 정책의 재분배 능력', '세금·사회보험·직업전환 지원·평생학습·공공 AI 접근권', ACCENT_GREEN),
    ('④ 교육의 전환 속도', '암기형 교육 → 문제정의·판단·시스템협업형 교육으로의 전환', ACCENT_GOLD),
    ('⑤ "일 없는 존엄" 인정', '일하지 않아도 인간의 존엄과 소득을 보장할 수 있는가? — 피할 수 없는 질문', ACCENT_PURPLE),
]
colors = [ACCENT_RED, ACCENT_CYAN, ACCENT_GREEN, ACCENT_GOLD, ACCENT_PURPLE]
for i, (title, desc, col) in enumerate(variables):
    y = Inches(1.8 + i * 1.0)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(10.3), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = col
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, title, 14, True, color=col)
    run2 = p.add_run()
    set_run(run2, '  ' + desc, 11, False, color=TEXT_DIM)

source_footer(s, '출처: Acemoglu & Johnson (2023), Piketty (2014)')


# ── SLIDE 62: REVISED POLARIZATION MAP ──
s = add_slide()
bridge_text(s, 'Part 2~8의 모든 분석을 종합한 최종 프레임', Inches(0.15))
part_marker(s, '최종 프레임', Inches(0.6))
title_text(s, '수정된 양극화 지도', Inches(1.0), 28)

headers = ['축', '표면', '심층', '결과']
rows = [
    ['💰 자산', '도구 민주화', '인프라 집중', '사용은 평등, 소유는 집중'],
    ['🧠 역량', '기초 평준화', '메타역량 희소화', '평균↑, 상위 초격차'],
    ['🔗 접속권', '기본 접속 대중화', '권한 깊이 계층화', '접근은 평등, 권한은 차등'],
    ['🎭 정체성', '다양화·유동화', '영향력 위계 재형성', '형태는 자유, 인정은 집중'],
]
add_simple_table(s, headers, rows, Inches(1.5), Inches(1.8), Inches(10.3),
                 [Inches(1.3), Inches(2.5), Inches(2.5), Inches(4)])

add_quote_box(s, '미래는\n**"기초 자원의 평준화"**와\n**"메타 권력의 집중"**이\n**동시에** 진행되는 사회다',
              Inches(3), Inches(4.5), Inches(7.3), Inches(2), ACCENT_RED, 20)


# ── SLIDE 63: FINAL COMPREHENSIVE SUMMARY ──
s = add_slide()
part_marker(s, '전체 강의 최종 정리', Inches(0.3))
bridge_text(s, '62페이지의 분석을 4줄로 압축하면', Inches(0.7))

summary = [
    ('🔧 기술사', '인간 기능의 대리화', ACCENT_RED),
    ('💼 일의 재정의', '"하는 것"에서 "정하는 것"으로', ACCENT_CYAN),
    ('👑 권력의 이동', '"의미의 독점"에서 "시스템의 독점"으로', ACCENT_PURPLE),
    ('⚖️ 양극화', '기초 평준화 + 메타 권력 집중', ACCENT_GOLD),
]
for i, (label, desc, col) in enumerate(summary):
    x = Inches(0.8 + i * 3.1)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(2.8), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = col
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, label, 15, True, color=col)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    set_run(run2, desc, 13, True, color=TEXT_MAIN)

# Key insight box
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x20, 0x0a, 0x12)
shape.line.color.rgb = ACCENT_RED
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(20)
tf.margin_top = Pt(14)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(26)
run = p.add_run()
set_run(run, 'AI 시대의 핵심 불평등은', 16, False, color=TEXT_SUB)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.line_spacing = Pt(26)
run2 = p2.add_run()
set_run(run2, '"누가 더 많이 가졌는가"보다', 16, False, color=TEXT_SUB)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.line_spacing = Pt(28)
run3 = p3.add_run()
set_run(run3, '"누가 목적함수와 예외와 책임을 결정하는가"에 있다', 18, True, color=ACCENT_RED)

# Sources
source_text = ('주요 참고: Eisenstein (1979), Mokyr (1990), McLuhan (1964), Weber (1922), Habermas (1962), '
               'Tilly (1992), Piketty (2014), Castells (2010), Zuboff (2019), Lessig (2006), '
               'Brynjolfsson & McAfee (2014), Rifkin (2014), Acemoglu & Johnson (2023), '
               'Autor (2024), Susskind (2020), Frey & Osborne (2017), Russell & Norvig (2020)\n'
               '데이터: Our World in Data, ITU, Statista, Stanford HAI (2024), Goldman Sachs (2024), McKinsey (2023), WEF (2023)')

tf = add_textbox(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5))
p = tf.paragraphs[0]
p.line_spacing = Pt(16)
run = p.add_run()
set_run(run, source_text, 8, False, color=TEXT_FAINT)


# ══════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════
output_path = r'C:\SIY_DEV\Lectures\tech-evolution-slides.pptx'
prs.save(output_path)
print(f'✅ PPTX saved to: {output_path}')
print(f'   Total slides: {len(prs.slides)}')
